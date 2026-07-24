"""Standalone downloads that don't need the full report bundle: just the
recommendations table (mapped to process steps, with problem statements),
and just the golden/reference benchmark dataset used by the KPI engine.
"""
from __future__ import annotations

import io

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.evaluation.kpi_engine import golden_benchmark_rows
from app.reports.report_data import ReportContext
from app.schemas.enums import efficiency_plan_for_category, main_category_for_category, sub_category_label

BRAND_BLUE = RGBColor(0x1D, 0x4E, 0xD8)
HEADER_FMT = {"bold": True, "bg_color": "#1D4ED8", "font_color": "white", "border": 1}


def _write_sheet(writer, header_format, df: pd.DataFrame, sheet_name: str):
    df.to_excel(writer, sheet_name=sheet_name, index=False, startrow=1, header=False)
    worksheet = writer.sheets[sheet_name]
    for col_idx, col_name in enumerate(df.columns):
        worksheet.write(0, col_idx, col_name, header_format)
        max_len = max([len(str(col_name))] + [len(str(v)) for v in df[col_name].astype(str).tolist()[:200]])
        worksheet.set_column(col_idx, col_idx, min(max_len + 2, 60))


def _step_name_lookup(ctx: ReportContext) -> dict[int, str]:
    return {d.step_number: d.step_name for d in ctx.diagnostics}


def generate_recommendations_excel(ctx: ReportContext) -> bytes:
    step_names = _step_name_lookup(ctx)
    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine="xlsxwriter") as writer:
        workbook = writer.book
        header_format = workbook.add_format(HEADER_FMT)

        rec_df = pd.DataFrame(
            [
                {
                    "Step": r.step_number or "Process-level",
                    "Step Name": step_names.get(r.step_number, "") if r.step_number else "Process-level",
                    "Main Category": main_category_for_category(r.category),
                    "Sub-Category": sub_category_label(r.category),
                    "Title": r.title,
                    "Problem Statement": r.problem_statement,
                    "Description": r.description,
                    "Rationale": r.rationale,
                    "Efficiency Plan": efficiency_plan_for_category(r.category),
                    "Proposed By": r.proposed_by_agent,
                    "Horizon": r.roadmap_horizon.value,
                    "Business Impact": r.prioritization.business_impact,
                    "Effort": r.prioritization.implementation_effort,
                    "ROI": r.prioritization.roi,
                    "Quadrant": r.prioritization.quadrant,
                    "FTE Savings": r.savings.fte_savings,
                    "Annual Savings ($)": r.savings.annual_cost_savings,
                    "Confidence": r.confidence_score,
                    "Source": r.source_type.value,
                    "Possible Duplicate": r.is_duplicate,
                }
                for r in ctx.recommendations
            ]
        )
        _write_sheet(writer, header_format, rec_df, "Recommendations")

    return buffer.getvalue()


def generate_recommendations_ppt(ctx: ReportContext) -> bytes:
    step_names = _step_name_lookup(ctx)
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Recommendations by Process Step"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = BRAND_BLUE
    slide.placeholders[1].text = f"{ctx.metadata.process_name}\n{ctx.metadata.team_name} | {ctx.metadata.lob}"

    active = [r for r in ctx.recommendations if not r.is_duplicate]
    by_step: dict[str, list] = {}
    for r in active:
        key = f"Step {r.step_number}: {step_names.get(r.step_number, '')}" if r.step_number else "Process-level"
        by_step.setdefault(key, []).append(r)

    for step_label, recs in by_step.items():
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = step_label
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = BRAND_BLUE

        header = ["Title", "Problem Statement", "Category", "Horizon"]
        rows = [[r.title[:40], (r.problem_statement or "")[:60], sub_category_label(r.category), r.roadmap_horizon.value] for r in recs[:8]]
        n_rows = len(rows) + 1
        table_shape = slide.shapes.add_table(n_rows, len(header), Inches(0.4), Inches(1.3), Inches(9.2), Inches(0.4 * n_rows))
        table = table_shape.table
        for c, h in enumerate(header):
            table.cell(0, c).text = h
            table.cell(0, c).text_frame.paragraphs[0].font.bold = True
            table.cell(0, c).text_frame.paragraphs[0].font.size = Pt(12)
        for r_idx, row in enumerate(rows, start=1):
            for c_idx, val in enumerate(row):
                cell = table.cell(r_idx, c_idx)
                cell.text = str(val)
                cell.text_frame.paragraphs[0].font.size = Pt(10)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


def generate_golden_dataset_excel() -> bytes:
    rows = golden_benchmark_rows()
    df = pd.DataFrame(rows)
    df.columns = [c.replace("_", " ").title() for c in df.columns]
    buffer = io.BytesIO()
    with pd.ExcelWriter(buffer, engine="xlsxwriter") as writer:
        workbook = writer.book
        header_format = workbook.add_format(HEADER_FMT)
        _write_sheet(writer, header_format, df, "Golden Benchmark Dataset")
    return buffer.getvalue()


def generate_golden_dataset_ppt() -> bytes:
    rows = golden_benchmark_rows()
    prs = Presentation()

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Golden Benchmark Dataset"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = BRAND_BLUE
    slide.placeholders[1].text = "Reference KPI targets used to benchmark each diagnostic"

    header = ["Category", "PCE %", "Automation %", "AI Readiness %", "Touch-Time %", "First-Pass Yield %", "Rework %"]
    keys = [
        "category", "process_cycle_efficiency_pct", "automation_coverage_pct", "ai_readiness_pct",
        "touch_time_ratio_pct", "first_pass_yield_pct", "rework_rate_pct",
    ]
    body = [[r[k] for k in keys] for r in rows]

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Reference Benchmarks by Category"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = BRAND_BLUE
    n_rows = len(body) + 1
    table_shape = slide.shapes.add_table(n_rows, len(header), Inches(0.3), Inches(1.3), Inches(9.4), Inches(0.4 * n_rows))
    table = table_shape.table
    for c, h in enumerate(header):
        table.cell(0, c).text = h
        table.cell(0, c).text_frame.paragraphs[0].font.bold = True
        table.cell(0, c).text_frame.paragraphs[0].font.size = Pt(11)
    for r_idx, row in enumerate(body, start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(10)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
