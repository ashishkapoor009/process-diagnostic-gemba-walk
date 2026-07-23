"""PowerPoint deck generation via python-pptx: a consulting-style
steering-committee deck (title, executive summary, key findings, quick
wins, roadmap, savings).
"""
from __future__ import annotations

import io

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.reports.report_data import ReportContext

BRAND_BLUE = RGBColor(0x1D, 0x4E, 0xD8)
DARK = RGBColor(0x0F, 0x17, 0x2A)


def _title_slide(prs: Presentation, ctx: ReportContext):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Process Diagnostic / Gemba Walk"
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = BRAND_BLUE
    subtitle = slide.placeholders[1]
    subtitle.text = f"{ctx.metadata.process_name}\n{ctx.metadata.team_name} | {ctx.metadata.lob}"


def _bullet_slide(prs: Presentation, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = BRAND_BLUE
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets or ["No content available."]):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
    return slide


def _table_slide(prs: Presentation, title: str, header: list[str], rows: list[list[str]]):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = BRAND_BLUE

    n_rows = min(len(rows), 8) + 1
    n_cols = len(header)
    table_shape = slide.shapes.add_table(n_rows, n_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.4 * n_rows))
    table = table_shape.table
    for c, h in enumerate(header):
        table.cell(0, c).text = h
        table.cell(0, c).text_frame.paragraphs[0].font.bold = True
        table.cell(0, c).text_frame.paragraphs[0].font.size = Pt(12)

    for r_idx, row in enumerate(rows[:8], start=1):
        for c_idx, val in enumerate(row):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
    return slide


def generate_ppt_report(ctx: ReportContext) -> bytes:
    prs = Presentation()

    _title_slide(prs, ctx)

    summary_paragraphs = [p.strip() for p in ctx.executive_summary.split("\n\n") if p.strip()]
    _bullet_slide(prs, "Executive Summary", summary_paragraphs[:6])

    s = ctx.savings_summary
    _bullet_slide(
        prs, "Key Metrics & Expected Impact",
        [
            f"Current State: {ctx.metadata.current_fte} FTE | {ctx.metadata.current_volume} volume/period | {ctx.metadata.aht_minutes} min AHT",
            f"Recommendations Identified: {s.get('total_recommendations', 0)} ({s.get('quick_win_count', 0)} Quick Wins, {s.get('strategic_count', 0)} Strategic)",
            f"Estimated FTE Savings (Released): {s.get('total_fte_savings', 0)}",
            f"In-Year Savings: ${s.get('in_year_savings', 0):,.0f}",
            f"12-Month Savings (Full Run-Rate): ${s.get('twelve_month_savings', 0):,.0f}",
            f"Blended Efficiency Improvement: {s.get('blended_efficiency_improvement_pct', 0)}% (target {s.get('target_efficiency_range_pct', '25-30%')})",
        ],
    )

    top_wastes = {}
    for d in ctx.diagnostics:
        for w in d.lean_wastes:
            top_wastes[w.value] = top_wastes.get(w.value, 0) + 1
    waste_bullets = [f"{k}: {v} step(s)" for k, v in sorted(top_wastes.items(), key=lambda kv: -kv[1])[:8]]
    _bullet_slide(prs, "Lean Waste Findings (TIMWOODS)", waste_bullets)

    _table_slide(
        prs, "Quick Wins", ["Title", "Category", "Impact", "Effort"],
        [[r.title[:35], r.category.value, r.prioritization.business_impact, r.prioritization.implementation_effort]
         for r in ctx.quick_wins],
    )

    for horizon, recs in list(ctx.roadmap_grouped().items())[:6]:
        _bullet_slide(prs, f"Roadmap: {horizon}", [f"{r.title} ({r.category.value})" for r in recs[:8]])

    _bullet_slide(
        prs, "Next Steps",
        [
            "Validate quick-win recommendations with process owners (Week 1-2)",
            "Kick off 30/60/90-day implementation sprints per roadmap",
            "Stand up governance/RACI for automation and AI initiatives",
            "Track realized savings against the estimates in this report",
        ],
    )

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
