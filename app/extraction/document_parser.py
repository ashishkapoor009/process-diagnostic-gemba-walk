"""Multi-format document parsing: PDF, DOCX, PPTX, CSV/Excel, images,
BPMN/Visio-exported XML, and a generic Unstructured.io fallback for
anything else. Every parser normalizes its output to plain text (+ any
tables found) so the LLM-based step extractor can work off one shape.
"""
from __future__ import annotations

from dataclasses import dataclass, field
from pathlib import Path

import pandas as pd
import pdfplumber
import pymupdf  # PyMuPDF, imported as `fitz` upstream but pymupdf works directly
from docx import Document as DocxDocument
from pptx import Presentation

from app.extraction.ocr import ocr_image, ocr_pil_image
from app.utils.logging import get_logger

logger = get_logger(__name__)

SUPPORTED_EXTENSIONS = {
    ".pdf", ".docx", ".doc", ".pptx", ".ppt", ".png", ".jpg", ".jpeg",
    ".csv", ".xlsx", ".xls", ".bpmn", ".xml", ".vsdx", ".txt",
}


@dataclass
class ExtractedDocument:
    filename: str
    file_type: str
    raw_text: str = ""
    tables: list[pd.DataFrame] = field(default_factory=list)
    used_ocr: bool = False

    @property
    def combined_text(self) -> str:
        parts = [self.raw_text]
        for i, table in enumerate(self.tables):
            parts.append(f"\n[Table {i + 1}]\n{table.to_markdown(index=False)}")
        return "\n".join(p for p in parts if p).strip()


def parse_pdf(path: Path) -> ExtractedDocument:
    """Extract text with PyMuPDF (fast, handles most native PDFs) and tables
    with pdfplumber. If a page yields near-zero text, treat it as a scanned
    image and fall back to OCR via PyMuPDF's page render.
    """
    text_chunks: list[str] = []
    tables: list[pd.DataFrame] = []
    used_ocr = False

    with pymupdf.open(path) as doc:
        for page_num, page in enumerate(doc):
            page_text = page.get_text().strip()
            if len(page_text) < 20:
                # Likely a scanned/image-only page -> OCR it.
                pix = page.get_pixmap(dpi=200)
                import io

                from PIL import Image

                img = Image.open(io.BytesIO(pix.tobytes("png")))
                ocr_text = ocr_pil_image(img)
                if ocr_text:
                    used_ocr = True
                    page_text = ocr_text
            text_chunks.append(f"[Page {page_num + 1}]\n{page_text}")

    try:
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                for table in page.extract_tables():
                    if table and len(table) > 1:
                        df = pd.DataFrame(table[1:], columns=table[0])
                        tables.append(df)
    except Exception as exc:  # pragma: no cover
        logger.warning(f"pdfplumber table extraction failed for {path.name}: {exc}")

    return ExtractedDocument(
        filename=path.name, file_type="pdf", raw_text="\n\n".join(text_chunks),
        tables=tables, used_ocr=used_ocr,
    )


def parse_docx(path: Path) -> ExtractedDocument:
    doc = DocxDocument(str(path))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    tables: list[pd.DataFrame] = []
    for table in doc.tables:
        rows = [[cell.text for cell in row.cells] for row in table.rows]
        if rows and len(rows) > 1:
            tables.append(pd.DataFrame(rows[1:], columns=rows[0]))
    return ExtractedDocument(
        filename=path.name, file_type="docx", raw_text="\n".join(paragraphs), tables=tables,
    )


def parse_pptx(path: Path) -> ExtractedDocument:
    prs = Presentation(str(path))
    chunks: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        slide_lines = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                slide_lines.append(shape.text_frame.text.strip())
            if shape.has_table:
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                slide_lines.append("\n".join(" | ".join(r) for r in rows))
        chunks.append("\n".join(slide_lines))
    return ExtractedDocument(filename=path.name, file_type="pptx", raw_text="\n\n".join(chunks))


def parse_image(path: Path) -> ExtractedDocument:
    text = ocr_image(path)
    return ExtractedDocument(filename=path.name, file_type="image", raw_text=text, used_ocr=True)


def parse_spreadsheet(path: Path) -> ExtractedDocument:
    if path.suffix.lower() == ".csv":
        df = pd.read_csv(path)
        tables = [df]
    else:
        sheets = pd.read_excel(path, sheet_name=None)
        tables = list(sheets.values())
    return ExtractedDocument(filename=path.name, file_type="spreadsheet", raw_text="", tables=tables)


def parse_bpmn_or_xml(path: Path) -> ExtractedDocument:
    """BPMN files (and Visio's flat XML export) are XML with human-readable
    labels on <bpmn:task name="..."> / <bpmn:sequenceFlow> style elements.
    We don't attempt full BPMN schema parsing - we extract every `name`
    attribute in document order, which reliably captures step labels for
    both BPMN 2.0 XML and Visio XML exports.
    """
    import re

    raw = path.read_text(encoding="utf-8", errors="ignore")
    names = re.findall(r'name="([^"]+)"', raw)
    text = "\n".join(f"- {n}" for n in names if n.strip())
    return ExtractedDocument(filename=path.name, file_type="bpmn", raw_text=text)


def parse_txt(path: Path) -> ExtractedDocument:
    return ExtractedDocument(filename=path.name, file_type="txt", raw_text=path.read_text(encoding="utf-8", errors="ignore"))


def parse_generic_unstructured(path: Path) -> ExtractedDocument:
    """Fallback for formats not explicitly handled above (e.g. legacy .doc,
    .vsdx) using unstructured.io's auto-partitioning.
    """
    try:
        from unstructured.partition.auto import partition

        elements = partition(filename=str(path))
        text = "\n".join(str(el) for el in elements)
        return ExtractedDocument(filename=path.name, file_type=path.suffix.lstrip("."), raw_text=text)
    except Exception as exc:  # pragma: no cover - optional heavy dependency
        logger.warning(f"Unstructured fallback failed for {path.name}: {exc}")
        return ExtractedDocument(filename=path.name, file_type=path.suffix.lstrip("."), raw_text="")


_PARSERS = {
    ".pdf": parse_pdf,
    ".docx": parse_docx,
    ".pptx": parse_pptx,
    ".png": parse_image,
    ".jpg": parse_image,
    ".jpeg": parse_image,
    ".csv": parse_spreadsheet,
    ".xlsx": parse_spreadsheet,
    ".xls": parse_spreadsheet,
    ".bpmn": parse_bpmn_or_xml,
    ".xml": parse_bpmn_or_xml,
    ".txt": parse_txt,
}


def parse_document(path: str | Path) -> ExtractedDocument:
    path = Path(path)
    ext = path.suffix.lower()
    parser = _PARSERS.get(ext, parse_generic_unstructured)
    logger.info(f"Parsing '{path.name}' with {parser.__name__}")
    return parser(path)
